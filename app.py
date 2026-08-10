import streamlit as st
import pandas as pd
import numpy as np
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import plotly.graph_objects as go

# ===== FUNÇÃO DE FORMATAÇÃO BRASILEIRA (Ex: 350.658,00) =====
def fmt_br(valor):
    try:
        return f"{float(valor):,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    except (ValueError, TypeError):
        return "0,00"

# ===== CONFIGURAÇÃO DA PÁGINA =====
st.set_page_config(page_title="Apresentação Comercial - KAO", page_icon="🚚", layout="wide")

# ===== IDENTIDADE VISUAL (DARK MODE / SLIDE POWERPOINT) =====
BG_DARK = "#0F172A"       # Fundo da aplicação (Azul Escuro / Preto)
CARD_DARK = "#1E293B"     # Fundo dos Cards/Slides
BORDER_DARK = "#334155"   # Borda elegante
TEXT_MAIN = "#F8FAFC"     # Texto principal (Branco)
TEXT_MUTED = "#94A3B8"    # Texto secundário (Cinza Claro)
GREEN_NEON = "#22C55E"    # Verde vibrante nítido
NAVY_ACCENT = "#38BDF8"   # Azul claro nítido

st.markdown(f"""
<style>
/* Fundo Geral da Aplicação */
.stApp {{ background-color: {BG_DARK}; color: {TEXT_MAIN}; }}
.block-container {{ padding-top: 1.5rem; max-width: 1350px; }}

/* Cabeçalhos */
h1, h2, h3, h4, h5 {{ color: {TEXT_MAIN} !important; font-family: 'Segoe UI', Roboto, sans-serif; }}

/* Card Estilo Slide PowerPoint */
.header-slide {{
    background: linear-gradient(135deg, #0284C7 0%, #0F172A 100%);
    border-left: 8px solid {GREEN_NEON};
    color: white;
    padding: 22px 28px;
    border-radius: 12px;
    margin-bottom: 24px;
    box-shadow: 0 10px 25px rgba(0,0,0,0.5);
}}

.section-card {{
    background: {CARD_DARK};
    padding: 24px;
    border-radius: 14px;
    border: 1px solid {BORDER_DARK};
    margin-bottom: 24px;
    box-shadow: 0 8px 20px rgba(0, 0, 0, 0.4);
}}

/* Ajustes nos componentes Streamlit para Dark Mode */
.stTextInput label, .stSelectbox label, .stNumberInput label, .stDateInput label {{
    color: {TEXT_MAIN} !important;
    font-weight: 600;
}}

/* Forçar visualização nítida de tabelas/editores */
[data-testid="stDataEditor"] {{
    background-color: {CARD_DARK};
    border-radius: 8px;
}}
</style>
""", unsafe_allow_html=True)

# ===== HEADER PRINCIPAL (SLIDE CAPA) =====
st.markdown(f"""
<div class="header-slide">
    <div style="display:flex; justify-content:space-between; align-items:center;">
        <div>
            <h1 style="color:white !important; margin:0; font-size:1.9rem; font-weight:700;">📊 APRESENTAÇÃO — TIME COMERCIAL</h1>
            <p style="margin:6px 0 0; color:#E2E8F0; font-size:1.0rem;">Painel de Acompanhamento e Desempenho Operacional / Comercial</p>
        </div>
        <div style="text-align:right;">
            <span style="background:{GREEN_NEON}; color:#000; padding:8px 16px; border-radius:20px; font-weight:800; font-size:0.85rem; letter-spacing:0.5px;">SISTEMA KAO</span>
        </div>
    </div>
</div>
""", unsafe_allow_html=True)

# ===== ESTADO DA SESSÃO =====
months = ["Janeiro", "Fevereiro", "Março", "Abril", "Maio", "Junho", 
          "Julho", "Agosto", "Setembro", "Outubro", "Novembro", "Dezembro"]

if "df_historico" not in st.session_state:
    st.session_state.df_historico = pd.DataFrame({
        "Mês": months,
        "Meta Total": [350658.00] * 12,
        "Fat. Total": [0.0] * 12,
        "% Total": [0.0] * 12,
        "UP": [0.0] * 12,
        "LOSS": [0.0] * 12
    })

if "clientes" not in st.session_state:
    st.session_state.clientes = pd.DataFrame([
        {"CLIENTE": f"Cliente {i+1}", "RECEITA": 0.0, "MÊS ANTERIOR": 0.0, "ANO ANTERIOR": 0.0, "RENTABILIDADE": 0.0, "SLA": 0.0, "CONSIDERAÇÕES": ""}
        for i in range(5)
    ])

if "faturado_auto" not in st.session_state:
    st.session_state.faturado_auto = 0.0

if "fechados" not in st.session_state:
    st.session_state.fechados = pd.DataFrame([
        {"CLIENTE": "", "PROJEÇÃO MÊS": 0.0, "FATURADO MÊS": 0.0, "PRODUTO": "", "ORIGEM": "", "DESTINO": ""}
        for _ in range(5)
    ])

if "upcoming" not in st.session_state:
    st.session_state.upcoming = pd.DataFrame([
        {"CLIENTE": "", "PROJEÇÃO MÊS": 0.0, "DATA PREVISTA": "", "PRODUTO": "", "ORIGEM": "", "DESTINO": ""}
        for _ in range(5)
    ])

# ===== SEÇÃO 0: IMPORTAÇÃO =====
with st.expander("📂 Importar Dados de Planilha Excel (Carga Rápida)", expanded=False):
    uploaded_file = st.file_uploader("Arraste ou selecione a planilha (.xlsx)", type=["xlsx"])
    if uploaded_file:
        try:
            excel_data = pd.read_excel(uploaded_file, sheet_name=None)
            if "Export" in excel_data:
                df_export = excel_data["Export"]
                if "Razão Grupo" in df_export.columns:
                    df_valid = df_export[
                        df_export["Razão Grupo"].notna() & 
                        (~df_export["Razão Grupo"].astype(str).str.startswith("Total")) &
                        (~df_export["Razão Grupo"].astype(str).str.startswith("Filtros"))
                    ].copy()

                    df_valid["Mês atual"] = pd.to_numeric(df_valid["Mês atual"], errors='coerce').fillna(0.0)
                    df_valid["Mês anterior"] = pd.to_numeric(df_valid["Mês anterior"], errors='coerce').fillna(0.0)

                    top5 = df_valid.sort_values(by="Mês atual", ascending=False).head(5)

                    novos_clientes = []
                    for idx, row in top5.iterrows():
                        novos_clientes.append({
                            "CLIENTE": str(row["Razão Grupo"]),
                            "RECEITA": float(row["Mês atual"]),
                            "MÊS ANTERIOR": float(row["Mês anterior"]),
                            "ANO ANTERIOR": 0.0,
                            "RENTABILIDADE": 0.0,
                            "SLA": 0.0,
                            "CONSIDERAÇÕES": ""
                        })

                    while len(novos_clientes) < 5:
                        novos_clientes.append({
                            "CLIENTE": "", "RECEITA": 0.0, "MÊS ANTERIOR": 0.0,
                            "ANO ANTERIOR": 0.0, "RENTABILIDADE": 0.0, "SLA": 0.0, "CONSIDERAÇÕES": ""
                        })

                    st.session_state.clientes = pd.DataFrame(novos_clientes)
                    st.session_state.faturado_auto = float(df_valid["Mês atual"].sum())
                    st.success("✅ Dados importados com sucesso!")

            elif "Historico" in excel_data:
                st.session_state.df_historico = excel_data["Historico"]
                st.success("✅ Histórico carregado!")
        except Exception as e:
            st.error(f"Erro ao ler arquivo: {e}")

# ===== SEÇÃO 1: IDENTIFICAÇÃO =====
st.markdown('<div class="section-card">', unsafe_allow_html=True)
st.subheader("1. Identificação do Executivo")
c1, c2, c3 = st.columns([2, 1, 1])
with c1:
    nome = st.text_input("NOME DO EXECUTIVO / KAM", value="Cleyton Kauffmann")
with c2:
    mes = st.selectbox("MÊS DE REFERÊNCIA", months, index=4)
with c3:
    data_ref = st.date_input("DATA DE APRESENTAÇÃO")
st.markdown('</div>', unsafe_allow_html=True)

# ===== SEÇÃO 2: SLIDE META / FATURAMENTO (MÊS ATUAL) =====
st.markdown('<div class="section-card">', unsafe_allow_html=True)
st.subheader("2. META | FATURAMENTO — MÊS ATUAL")

col_meta, col_fat, col_proj = st.columns(3)

with col_meta:
    meta = st.number_input("📌 META DA GERÊNCIA (R$)", min_value=0.0, value=350658.00, step=1000.0)

with col_fat:
    faturado = st.number_input(
        "💰 FATURADO ALCANÇADO (R$)", 
        min_value=0.0, 
        value=st.session_state.get("faturado_auto", 0.0), 
        step=1000.0
    )

with col_proj:
    projecao = st.number_input("📈 PROJEÇÃO MÊS (R$)", min_value=0.0, value=faturado, step=1000.0)

pct_realizado = (faturado / meta * 100) if meta > 0 else 0.0
pct_projecao = (projecao / meta * 100) if meta > 0 else 0.0
gap_projecao = projecao - meta
delta_pct = pct_realizado - 100

st.divider()

# ===== ETIQUETAS E BOLA HIGH-CONTRAST (DARK STYLE) =====
m1, m2, m3 = st.columns(3)

cor_bola_bg = "linear-gradient(135deg, #166534 0%, #22C55E 100%)" if pct_realizado >= 100 else "linear-gradient(135deg, #991B1B 0%, #EF4444 100%)"
cor_badge_bg = "#14532D" if delta_pct >= 0 else "#7F1D1D"
cor_badge_txt = "#4ADE80" if delta_pct >= 0 else "#FCA5A5"

with m1:
    st.markdown(f"""
    <div style="text-align: center; background: #0F172A; padding: 20px; border-radius: 14px; border: 1px solid {BORDER_DARK};">
        <span style="font-size: 0.85rem; font-weight: 700; color: {TEXT_MUTED}; letter-spacing: 0.5px; text-transform: uppercase;">% ALCANÇADO (META x FAT)</span>
        <div style="
            width: 120px; 
            height: 120px; 
            background: {cor_bola_bg}; 
            border-radius: 50%; 
            display: flex; 
            align-items: center; 
            justify-content: center; 
            margin: 15px auto; 
            box-shadow: 0 0 20px rgba(34, 197, 94, 0.4);
        ">
            <span style="color: #FFFFFF; font-size: 1.6rem; font-weight: 800;">{pct_realizado:.1f}%</span>
        </div>
        <span style="
            background-color: {cor_badge_bg}; 
            color: {cor_badge_txt}; 
            padding: 6px 14px; 
            border-radius: 20px; 
            font-size: 0.85rem; 
            font-weight: 700;
        ">
            {"↑" if delta_pct >= 0 else "↓"} {delta_pct:+.1f}% vs Meta
        </span>
    </div>
    """.replace(".", ","), unsafe_allow_html=True)

with m2:
    st.markdown(f"""
    <div style="text-align: center; background: #0F172A; padding: 20px; border-radius: 14px; border: 1px solid {BORDER_DARK};">
        <span style="font-size: 0.85rem; font-weight: 700; color: {TEXT_MUTED}; letter-spacing: 0.5px; text-transform: uppercase;">PROJEÇÃO DA META %</span>
        <div style="
            width: 120px; 
            height: 120px; 
            background: linear-gradient(135deg, #0369A1 0%, #38BDF8 100%); 
            border-radius: 50%; 
            display: flex; 
            align-items: center; 
            justify-content: center; 
            margin: 15px auto; 
            box-shadow: 0 0 20px rgba(56, 189, 248, 0.3);
        ">
            <span style="color: #FFFFFF; font-size: 1.6rem; font-weight: 800;">{pct_projecao:.1f}%</span>
        </div>
        <span style="color: {TEXT_MUTED}; font-size: 0.85rem; font-weight: 600;">Projeção de Fechamento</span>
    </div>
    """.replace(".", ","), unsafe_allow_html=True)

with m3:
    cor_gap = "#4ADE80" if gap_projecao >= 0 else "#FCA5A5"
    st.markdown(f"""
    <div style="text-align: center; background: #0F172A; padding: 20px; border-radius: 14px; border: 1px solid {BORDER_DARK};">
        <span style="font-size: 0.85rem; font-weight: 700; color: {TEXT_MUTED}; letter-spacing: 0.5px; text-transform: uppercase;">GAP (PROJEÇÃO x META)</span>
        <div style="height: 120px; display: flex; align-items: center; justify-content: center; margin: 15px auto;">
            <span style="color: #FFFFFF; font-size: 2.0rem; font-weight: 800;">R$ {fmt_br(gap_projecao)}</span>
        </div>
        <span style="
            background-color: {cor_badge_bg}; 
            color: {cor_gap}; 
            padding: 6px 14px; 
            border-radius: 20px; 
            font-size: 0.85rem; 
            font-weight: 700;
        ">
            {"↑" if gap_projecao >= 0 else "↓"} R$ {fmt_br(gap_projecao)}
        </span>
    </div>
    """, unsafe_allow_html=True)

st.markdown('</div>', unsafe_allow_html=True)

# ===== SEÇÃO 3: EVOLUÇÃO HISTÓRICA =====
st.markdown('<div class="section-card">', unsafe_allow_html=True)
st.subheader("3. META | FATURAMENTO — EVOLUÇÃO HISTÓRICA")

st.markdown("##### ⚡ Ações Rápidas de Preenchimento")

c_auto1, c_auto2, c_auto3 = st.columns([2, 2, 2])

# ----- COLUNA 1: LANÇAR META MÊS A MÊS -----
with c_auto1:
    mes_meta_lanc = st.selectbox("Lançar Meta no Mês", months, index=months.index(mes), key="sel_mes_meta")
    valor_meta_lanc = st.number_input("Valor da Meta (R$)", min_value=0.0, value=float(meta), step=1000.0, key="input_meta_lanc")
    
    col_btn_meta1, col_btn_meta2 = st.columns(2)
    with col_btn_meta1:
        if st.button("💾 Salvar Meta Mês", use_container_width=True):
            idx = st.session_state.df_historico[st.session_state.df_historico["Mês"] == mes_meta_lanc].index
            if len(idx) > 0:
                st.session_state.df_historico.loc[idx[0], "Meta Total"] = valor_meta_lanc
                st.success(f"Meta de {mes_meta_lanc} atualizada!")
                st.rerun()
    with col_btn_meta2:
        if st.button("🔄 Replicar p/ Todos", use_container_width=True, help="Aplica este valor de Meta para todos os 12 meses"):
            st.session_state.df_historico["Meta Total"] = valor_meta_lanc
            st.success("Meta replicada para todos os meses!")
            st.rerun()

# ----- COLUNA 2: LANÇAR FATURADO MÊS A MÊS -----
with c_auto2:
    mes_lancamento = st.selectbox("Lançar Faturado no Mês", months, index=months.index(mes), key="sel_mes_lanc")
    valor_fat_lanc = st.number_input("Valor Faturado (R$)", min_value=0.0, value=float(faturado), step=1000.0, key="input_fat_lanc")
    if st.button("💾 Salvar Faturado Mês", use_container_width=True):
        idx = st.session_state.df_historico[st.session_state.df_historico["Mês"] == mes_lancamento].index
        if len(idx) > 0:
            st.session_state.df_historico.loc[idx[0], "Fat. Total"] = valor_fat_lanc
            st.success(f"Faturado de {mes_lancamento} atualizado para R$ {fmt_br(valor_fat_lanc)}!")
            st.rerun()

# ----- COLUNA 3: AÇÕES GLOBAIS -----
with c_auto3:
    st.write("##### 🧹 Limpeza")
    st.write("")
    if st.button("🧹 Zerar Faturamentos", use_container_width=True):
        st.session_state.df_historico["Fat. Total"] = 0.0
        st.rerun()

st.divider()

# Cálculo automático
df_calc = st.session_state.df_historico.copy()
df_calc["% Total"] = np.where(df_calc["Meta Total"] > 0, (df_calc["Fat. Total"] / df_calc["Meta Total"]) * 100, 0.0)
df_calc["UP"] = np.maximum(df_calc["Fat. Total"] - df_calc["Meta Total"], 0)
df_calc["LOSS"] = np.maximum(df_calc["Meta Total"] - df_calc["Fat. Total"], 0)
st.session_state.df_historico = df_calc

st.markdown("##### 📋 Tabela de Evolução Consolidada")

st.session_state.df_historico = st.data_editor(
    st.session_state.df_historico,
    use_container_width=True,
    hide_index=True,
    column_config={
        "Mês": st.column_config.TextColumn(disabled=True),
        "Meta Total": st.column_config.NumberColumn("Meta Total (R$)", format="R$ %,.2f"),
        "Fat. Total": st.column_config.NumberColumn("Fat. Total (R$)", format="R$ %,.2f"),
        "% Total": st.column_config.NumberColumn("% Atingido", format="%.1f%%", disabled=True),
        "UP": st.column_config.NumberColumn("UP (R$)", format="R$ %,.2f", disabled=True),
        "LOSS": st.column_config.NumberColumn("LOSS (R$)", format="R$ %,.2f", disabled=True),
    },
    key="editor_historico_v10"
)

st.markdown("#### Evolução Mensal (Gráfico Nítido com Tela Cheia / Zoom)")

# ===== GRÁFICO PLOTLY - DARK MODE NÍTIDO =====
fig = go.Figure()

# Barra Fina Verde Neón
fig.add_trace(go.Bar(
    x=st.session_state.df_historico["Mês"],
    y=st.session_state.df_historico["Fat. Total"],
    name="Faturado Alcançado",
    marker_color=GREEN_NEON,
    width=0.35,
    hovertemplate="Mês: %{x}<br>Alcançado: R$ %{y:,.2f}<extra></extra>"
))

# Linha de Tendência Azul Cyan
fig.add_trace(go.Scatter(
    x=st.session_state.df_historico["Mês"],
    y=st.session_state.df_historico["Fat. Total"],
    name="Tendência Faturado",
    mode="lines+markers",
    line=dict(color=NAVY_ACCENT, dash="dash", width=3),
    marker=dict(size=8, color=NAVY_ACCENT, symbol="circle"),
    hovertemplate="Mês: %{x}<br>Tendência: R$ %{y:,.2f}<extra></extra>"
))

fig.update_layout(
    paper_bgcolor="rgba(0,0,0,0)",
    plot_bgcolor="#0F172A",
    font=dict(color="#F8FAFC", family="Segoe UI"),
    xaxis=dict(
        categoryorder="array", 
        categoryarray=months,
        gridcolor="#334155",
        zerolinecolor="#334155",
        tickfont=dict(color="#F8FAFC", size=12)
    ),
    yaxis=dict(
        title="Valor (R$)", 
        tickprefix="R$ ",
        gridcolor="#334155",
        zerolinecolor="#334155",
        tickfont=dict(color="#F8FAFC", size=12)
    ),
    legend=dict(
        orientation="h", 
        yanchor="bottom", 
        y=1.02, 
        xanchor="right", 
        x=1,
        font=dict(color="#F8FAFC", size=12)
    ),
    margin=dict(l=20, r=20, t=40, b=20),
    hovermode="x unified"
)

st.plotly_chart(fig, use_container_width=True)
st.markdown('</div>', unsafe_allow_html=True)

# ===== SEÇÃO 4: TOP 5 CLIENTES =====
st.markdown('<div class="section-card">', unsafe_allow_html=True)
st.subheader("4. PRINCIPAIS CLIENTES (TOP 5)")

st.session_state.clientes = st.data_editor(
    st.session_state.clientes,
    num_rows="fixed",
    use_container_width=True,
    hide_index=True,
    key="editor_clientes",
    column_config={
        "RECEITA": st.column_config.NumberColumn("RECEITA (R$)", format="R$ %,.2f"),
        "MÊS ANTERIOR": st.column_config.NumberColumn("MÊS ANTERIOR (R$)", format="R$ %,.2f"),
        "ANO ANTERIOR": st.column_config.NumberColumn("ANO ANTERIOR (R$)", format="R$ %,.2f"),
        "RENTABILIDADE": st.column_config.NumberColumn("RENTABILIDADE %", format="%.1f%%"),
        "SLA": st.column_config.NumberColumn("SLA %", format="%.1f%%"),
        "CONSIDERAÇÕES": st.column_config.TextColumn("CONSIDERAÇÕES / OBS")
    }
)

df_valid_cli = st.session_state.clientes[st.session_state.clientes["CLIENTE"].astype(str).str.strip() != ""]
tot_top5 = df_valid_cli["RECEITA"].sum()
tot_carteira = st.session_state.get("faturado_auto", tot_top5)

c_tot1, c_tot2 = st.columns(2)
with c_tot1:
    st.metric("TOTAL TOP 5 CLIENTES", f"R$ {fmt_br(tot_top5)}")
with c_tot2:
    st.metric("TOTAL CARTEIRA GERAL", f"R$ {fmt_br(tot_carteira)}")
st.markdown('</div>', unsafe_allow_html=True)

# ===== SEÇÃO 5: PIPELINE =====
st.markdown('<div class="section-card">', unsafe_allow_html=True)
st.subheader("5. PIPELINE — NOVOS NEGÓCIOS FECHADOS (MÊS)")

st.session_state.fechados = st.data_editor(
    st.session_state.fechados,
    num_rows="dynamic",
    use_container_width=True,
    hide_index=True,
    key="editor_fechados",
    column_config={
        "PROJEÇÃO MÊS": st.column_config.NumberColumn(format="R$ %,.2f"),
        "FATURADO MÊS": st.column_config.NumberColumn(format="R$ %,.2f"),
    }
)

tot_proj_fechados = pd.to_numeric(st.session_state.fechados["PROJEÇÃO MÊS"], errors='coerce').fillna(0.0).sum()
tot_fat_fechados = pd.to_numeric(st.session_state.fechados["FATURADO MÊS"], errors='coerce').fillna(0.0).sum()

cp_tot1, cp_tot2 = st.columns(2)
with cp_tot1:
    st.metric("TOTAL PROJEÇÃO MÊS", f"R$ {fmt_br(tot_proj_fechados)}")
with cp_tot2:
    st.metric("TOTAL FATURADO MÊS", f"R$ {fmt_br(tot_fat_fechados)}")

st.divider()

st.subheader("UPCOMING — FECHADO PARA INICIAR")
st.session_state.upcoming = st.data_editor(
    st.session_state.upcoming,
    num_rows="dynamic",
    use_container_width=True,
    hide_index=True,
    key="editor_upcoming",
    column_config={
        "PROJEÇÃO MÊS": st.column_config.NumberColumn(format="R$ %,.2f"),
    }
)

tot_proj_upcoming = pd.to_numeric(st.session_state.upcoming["PROJEÇÃO MÊS"], errors='coerce').fillna(0.0).sum()

cu_tot1, _ = st.columns([1, 1])
with cu_tot1:
    st.metric("TOTAL PROJEÇÃO MÊS (UPCOMING)", f"R$ {fmt_br(tot_proj_upcoming)}")

st.markdown('</div>', unsafe_allow_html=True)

# ===== SEÇÃO 6: EXPORTAÇÃO =====
st.markdown('<div class="section-card">', unsafe_allow_html=True)
st.subheader("6. EXPORTAR RELATÓRIO")

col_exp1, col_exp2 = st.columns(2)

with col_exp1:
    if st.button("📥 Gerar Excel Completo", use_container_width=True):
        output_xl = BytesIO()
        with pd.ExcelWriter(output_xl, engine="openpyxl") as writer:
            pd.DataFrame([{
                "Nome": nome, "Mês": mes, "Data": data_ref,
                "Meta Gerência": fmt_br(meta), "Faturado Alcançado": fmt_br(faturado),
                "% Alcançado": f"{pct_realizado:.1f}%".replace(".", ","), "Projeção": fmt_br(projecao)
            }]).to_excel(writer, sheet_name="Resumo_Executivo", index=False)
            
            calc_hist_fmt = st.session_state.df_historico.copy()
            for col in ["Meta Total", "Fat. Total", "UP", "LOSS"]:
                calc_hist_fmt[col] = calc_hist_fmt[col].apply(fmt_br)
            calc_hist_fmt.to_excel(writer, sheet_name="Historico", index=False)
            
            st.session_state.clientes.to_excel(writer, sheet_name="Principais_Clientes", index=False)
            st.session_state.fechados.to_excel(writer, sheet_name="Novos_Negocios", index=False)
            st.session_state.upcoming.to_excel(writer, sheet_name="Upcoming", index=False)
        
        st.download_button(
            label="⬇️ Baixar Planilha (.xlsx)",
            data=output_xl.getvalue(),
            file_name=f"painel_comercial_{mes.lower()}_{nome.replace(' ','_')}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True
        )

with col_exp2:
    if st.button("📊 Gerar Apresentação PowerPoint (.pptx)", use_container_width=True):
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        
        slide1 = prs.slides.add_slide(blank_slide_layout)
        txBox = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tf = txBox.text_frame
        tf.text = f"APRESENTAÇÃO TIME COMERCIAL - {mes.upper()}"
        
        p2 = tf.add_paragraph()
        p2.text = f"Executivo: {nome} | Faturado: R$ {fmt_br(faturado)} / Meta: R$ {fmt_br(meta)} ({pct_realizado:.1f}%)".replace(".", ",")
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(34, 197, 94)
        
        output_ppt = BytesIO()
        prs.save(output_ppt)
        
        st.download_button(
            label="⬇️ Baixar Apresentação (.pptx)",
            data=output_ppt.getvalue(),
            file_name=f"apresentacao_comercial_{mes.lower()}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )

st.markdown('</div>', unsafe_allow_html=True)
